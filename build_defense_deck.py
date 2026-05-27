"""Build the diploma-defense pptx for the ReconX project.

Mirrors the structure and visual feel of `Presetation fin.pptx`:
    16:9, dark-navy + bright-blue palette, Calibri Light, badge boxes,
    comparison tables, big-stat callouts.

Run:
    python3 build_defense_deck.py
Output:
    ReconX_Defense_Mavlanov.pptx (next to this script)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# ── Palette (matches the example deck) ────────────────────────────────────────
NAVY        = RGBColor(0x1A, 0x2B, 0x5F)
NAVY_DEEP   = RGBColor(0x0E, 0x1B, 0x44)
INDIGO      = RGBColor(0x1E, 0x3A, 0x8A)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
BLUE_BRIGHT = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT  = RGBColor(0x93, 0xC5, 0xFD)
BLUE_PALE   = RGBColor(0xEF, 0xF6, 0xFF)
SLATE       = RGBColor(0x33, 0x41, 0x55)
SLATE_LIGHT = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT  = RGBColor(0xCB, 0xD5, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xDC, 0x26, 0x26)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)

FOOTER = "Mavlanov Shakhrukh  |  CS-2304  |  Astana IT University  |  June 2026"
SUPERVISOR = "PhD Amirova Akzhibek"


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
             corner_radius=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if corner_radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        # set corner radius (adj value: fraction of half the smaller dim)
        shp.adjustments[0] = corner_radius
    shp.text_frame.margin_left = Pt(0)
    shp.text_frame.margin_right = Pt(0)
    shp.text_frame.margin_top = Pt(0)
    shp.text_frame.margin_bottom = Pt(0)
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri Light"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_footer(slide):
    add_text(slide, Inches(0.4), Inches(5.25), Inches(9.2), Inches(0.25),
             FOOTER, size=9, color=SLATE_LIGHT, align=PP_ALIGN.LEFT)


def add_slide_header(slide, title, subtitle):
    # accent stripe
    add_rect(slide, Inches(0.4), Inches(0.38), Inches(0.12), Inches(0.55),
             fill=BLUE)
    add_text(slide, Inches(0.65), Inches(0.30), Inches(9.0), Inches(0.55),
             title, size=26, bold=True, color=NAVY)
    add_text(slide, Inches(0.65), Inches(0.78), Inches(9.0), Inches(0.30),
             subtitle, size=12, color=SLATE_LIGHT, font="Calibri Light")


def chip(slide, x, y, w, h, text, fill, color=WHITE, size=10):
    """Pill-shaped badge."""
    shp = add_rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                   corner_radius=0.5)
    tf = shp.text_frame
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri Light"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


# ─────────────────────────────────────────────────────────────────────────────
# Build presentation
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation()
# 16:9 — match example (9144000 x 5143500 EMU = 10in x 5.625in)
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(5143500)

BLANK = prs.slide_layouts[6]  # blank layout


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY_DEEP)

# decorative accent bar
add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.35), fill=BLUE)

add_text(s, Inches(0.5), Inches(0.6), Inches(9), Inches(0.4),
         "ASTANA IT UNIVERSITY  |  2026", size=11, bold=True,
         color=BLUE_LIGHT, font="Calibri")

add_text(s, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6),
         "Development of a Method for",
         size=30, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(1.75), Inches(9), Inches(0.6),
         "External Penetration Testing of a Web",
         size=30, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(2.30), Inches(9), Inches(0.6),
         "Application for Vulnerabilities",
         size=30, bold=True, color=WHITE)

# subtitle / tag line
add_rect(s, Inches(0.5), Inches(3.05), Inches(0.06), Inches(0.30), fill=BLUE_BRIGHT)
add_text(s, Inches(0.7), Inches(3.02), Inches(9), Inches(0.35),
         "ReconX  •  40+ Reconnaissance Modules  •  AI-Augmented Reporting  •  Active Probes",
         size=12, color=BLUE_LIGHT)

# meta box
y = Inches(3.7)
labels = [
    ("Student:",    "Mavlanov Shakhrukh"),
    ("Group:",      "CS-2304"),
    ("Supervisor:", SUPERVISOR),
    ("Programme:",  "6B06301 — Cybersecurity"),
]
for i, (k, v) in enumerate(labels):
    col_x = Inches(0.5 + (i % 2) * 4.6)
    row_y = Inches(3.7 + (i // 2) * 0.55)
    add_text(s, col_x, row_y, Inches(1.4), Inches(0.30),
             k, size=11, color=BLUE_LIGHT)
    add_text(s, col_x + Inches(1.2), row_y, Inches(3.5), Inches(0.30),
             v, size=12, bold=True, color=WHITE)


# ── Slide 2 — Motivation & Problem Statement ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Motivation & Problem Statement",
                 "Why automated external pentesting matters today")
add_footer(s)

# "THE PROBLEM" KPI strip
add_text(s, Inches(0.4), Inches(1.10), Inches(2), Inches(0.3),
         "THE PROBLEM", size=10, bold=True, color=BLUE)

stats = [
    ("70%+",    "Web apps remain the #1 initial breach vector (Verizon DBIR 2024)"),
    ("Days",    "Average time to manually triage a moderate external surface"),
    ("Few",     "Open tools combine recon + probes + AI reporting end-to-end"),
    ("Noisy",   "Existing scanners drown analysts in false positives"),
]
for i, (big, lbl) in enumerate(stats):
    x = Inches(0.4 + i * 2.35)
    add_rect(s, x, Inches(1.45), Inches(2.20), Inches(1.40),
             fill=BLUE_PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             corner_radius=0.05)
    add_text(s, x + Inches(0.1), Inches(1.55), Inches(2.0), Inches(0.6),
             big, size=24, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    add_text(s, x + Inches(0.1), Inches(2.10), Inches(2.0), Inches(0.7),
             lbl, size=9.5, color=SLATE)

# GAPS in existing solutions
add_text(s, Inches(0.4), Inches(3.05), Inches(8), Inches(0.3),
         "GAPS IN EXISTING TOOLS", size=10, bold=True, color=RED)
gaps = [
    ("✗",  "Nuclei alone — fast, but no contextual triage and no AI summary"),
    ("✗",  "Burp Suite Pro — manual, paid, single-target focused"),
    ("✗",  "Many OSS scanners ship default rules with high false-positive rates"),
    ("✗",  "No single open project combines passive recon + active probes + per-module false-positive guards + AI-generated client report"),
]
for i, (mark, txt) in enumerate(gaps):
    y = Inches(3.35 + i * 0.35)
    add_text(s, Inches(0.40), y, Inches(0.4), Inches(0.30),
             mark, size=14, bold=True, color=RED)
    add_text(s, Inches(0.85), y, Inches(8.7), Inches(0.30),
             txt, size=11, color=SLATE)


# ── Slide 3 — Research Objectives & Tasks ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Research Objectives & Tasks",
                 "What this diploma project set out to achieve")
add_footer(s)

# Objective banner
add_rect(s, Inches(0.40), Inches(1.10), Inches(9.20), Inches(0.55),
         fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.10)
add_text(s, Inches(0.55), Inches(1.18), Inches(9.0), Inches(0.40),
         "Objective: Design and implement an automated external-pentesting framework that "
         "discovers attack surface, probes vulnerabilities, and produces an analyst-ready report.",
         size=11, color=WHITE)

# Numbered tasks
tasks = [
    ("1", "Surface Discovery",
     "Subdomain enumeration, DNS, port scan, tech stack detection, live-host triage."),
    ("2", "Modular Architecture",
     "40+ Python modules sharing a BaseModule contract: independent, configurable, parallel-safe."),
    ("3", "Active Vulnerability Probes",
     "XSS, SQLi, IDOR, SSRF, SSTI, JWT, OAuth, host-header injection, cache poisoning, prototype pollution."),
    ("4", "False-Positive Reduction",
     "Per-module verifiers: SPA-fallback detection, JSON-URI-echo filters, vendor-library skip, repo relevance scoring."),
    ("5", "AI-Augmented Reporting",
     "Ollama / OpenAI-compatible LLM consumes structured scan data and writes a client-ready Markdown report."),
]
for i, (n, h, d) in enumerate(tasks):
    y = Inches(1.85 + i * 0.62)
    # number bubble
    add_rect(s, Inches(0.45), y, Inches(0.50), Inches(0.50),
             fill=BLUE, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(0.45), y, Inches(0.50), Inches(0.50),
             n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.10), y - Inches(0.02), Inches(8.5), Inches(0.30),
             h, size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.10), y + Inches(0.25), Inches(8.5), Inches(0.35),
             d, size=10.5, color=SLATE)


# ── Slide 4 — ReconX Architecture Overview ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "ReconX Framework — Architecture Overview",
                 "Three orchestrated stages with fail-isolated modules")
add_footer(s)

# Banner
add_rect(s, Inches(0.40), Inches(1.10), Inches(9.20), Inches(0.40),
         fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.20)
add_text(s, Inches(0.55), Inches(1.16), Inches(9.0), Inches(0.30),
         "ReconX — Passive Recon  →  Active Probing  →  AI Synthesis",
         size=12, bold=True, color=WHITE)

# Three big stage cards
def stage(slide, x, n, title, body, tag, blocks):
    add_rect(slide, x, Inches(1.65), Inches(3.0), Inches(3.0),
             fill=BLUE_PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             corner_radius=0.05)
    add_rect(slide, x, Inches(1.65), Inches(3.0), Inches(0.45),
             fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.10)
    add_text(slide, x + Inches(0.20), Inches(1.71), Inches(2.5), Inches(0.35),
             f"L{n}  {title}", size=12, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.20), Inches(2.20), Inches(2.7), Inches(1.6),
             body, size=10.5, color=SLATE)
    add_rect(slide, x + Inches(0.20), Inches(3.85), Inches(2.6), Inches(0.30),
             fill=BLUE_BRIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             corner_radius=0.40)
    add_text(slide, x + Inches(0.20), Inches(3.85), Inches(2.6), Inches(0.30),
             tag, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.20), Inches(4.20), Inches(2.7), Inches(0.45),
             blocks, size=9.5, color=BLUE)

stage(s, Inches(0.40), 1, "Passive Reconnaissance",
      "Subdomain enumeration (amass, subfinder, crt.sh), DNS records, WHOIS, "
      "ASN lookups, TLS certificate inspection, ICMP + HTTP liveness probing.",
      "Surfaces every reachable host",
      "Modules: recon, webdetect, portscan, techstack, ssl_checker")

stage(s, Inches(3.50), 2, "Active Probing",
      "Per-vulnerability probes: XSS, SQLi, IDOR/BOLA, SSRF/SSTI, JWT, "
      "OAuth, host-header injection, cache poisoning, prototype pollution, "
      "JS-secret extraction, takeover detection.",
      "Generates concrete evidence",
      "Modules: fuzzer, xss, sql_injection, idor_probe, …")

stage(s, Inches(6.60), 3, "AI Synthesis & Reporting",
      "LLM (Ollama / OpenAI-compat) ingests structured findings; deterministic "
      "fallback if AI unavailable. Outputs HTML report, JSON dump, and "
      "Telegram summary.",
      "Analyst-ready in minutes",
      "Modules: ai_report, asset_risk, correlator, report")

add_text(s, Inches(0.4), Inches(4.85), Inches(9.2), Inches(0.30),
         "All 40+ modules share a single BaseModule contract  •  Independent failure domains  "
         "•  Parallel-safe  •  Config-driven scope and rate limits",
         size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── Slide 5 — Live Demo: Report.html overview + Top Risk Assets ─────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Live Demo — Generated HTML Report",
                 "Real output from a scan against miras.app")
add_footer(s)

import os
ASSETS = "/home/kali/Desktop/Diploma-pre-defend/reconx/build_deck_assets/screenshots"

# Left: full overview / KPI strip
if os.path.exists(f"{ASSETS}/06_html_overview.png"):
    s.shapes.add_picture(f"{ASSETS}/06_html_overview.png",
                         Inches(0.30), Inches(1.10), width=Inches(5.20),
                         height=Inches(3.30))
    add_text(s, Inches(0.30), Inches(4.45), Inches(5.20), Inches(0.30),
             "↑ Report header — risk grade, 20 KPI tiles (subdomains, "
             "live hosts, vulnerabilities, CVEs, findings).",
             size=9, color=SLATE_LIGHT, align=PP_ALIGN.LEFT)

# Right: Top Risk Assets table screenshot
if os.path.exists(f"{ASSETS}/01_top_risk_assets.png"):
    s.shapes.add_picture(f"{ASSETS}/01_top_risk_assets.png",
                         Inches(5.65), Inches(1.10), width=Inches(4.00),
                         height=Inches(2.30))
    add_text(s, Inches(5.65), Inches(3.45), Inches(4.00), Inches(0.30),
             "↑ Top Risk Assets table — per-host score, tier, "
             "open ports, CVE/EDB count, confirmed findings.",
             size=9, color=SLATE_LIGHT)

# Live Subdomains snippet
if os.path.exists(f"{ASSETS}/03_live_subdomains.png"):
    s.shapes.add_picture(f"{ASSETS}/03_live_subdomains.png",
                         Inches(5.65), Inches(3.85), width=Inches(4.00),
                         height=Inches(1.30))
    add_text(s, Inches(5.65), Inches(5.15), Inches(4.00), Inches(0.20),
             "↑ Live Subdomains — DNS / ping / HTTP liveness chips.",
             size=8.5, color=SLATE_LIGHT)


# ── Slide 5b — All Findings + AI + Telegram (more screenshots) ──────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY_DEEP)
add_text(s, Inches(0.5), Inches(0.40), Inches(9), Inches(0.5),
         "Live Demo — Findings, AI Report, Telegram", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.5), Inches(0.90), Inches(9), Inches(0.35),
         "Three artifacts every scan delivers", size=11, color=BLUE_LIGHT)

# Left big card: All Findings table screenshot
if os.path.exists(f"{ASSETS}/02_all_findings.png"):
    add_rect(s, Inches(0.30), Inches(1.30), Inches(5.50), Inches(3.55),
             fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    s.shapes.add_picture(f"{ASSETS}/02_all_findings.png",
                         Inches(0.40), Inches(1.40), width=Inches(5.30),
                         height=Inches(3.35))
add_text(s, Inches(0.30), Inches(4.90), Inches(5.50), Inches(0.20),
         "All Findings — module · severity · URL · confidence + filters",
         size=9, color=BLUE_LIGHT)

# Top-right: AI analysis screenshot
if os.path.exists(f"{ASSETS}/05_ai_analysis.png"):
    add_rect(s, Inches(5.95), Inches(1.30), Inches(3.75), Inches(1.85),
             fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    s.shapes.add_picture(f"{ASSETS}/05_ai_analysis.png",
                         Inches(6.00), Inches(1.35), width=Inches(3.65),
                         height=Inches(1.75))
add_text(s, Inches(5.95), Inches(3.18), Inches(3.75), Inches(0.20),
         "AI Security Analysis (Ollama deepseek-r1)",
         size=9, color=BLUE_LIGHT)

# Bottom-right: Telegram
if os.path.exists(f"{ASSETS}/04_telegram_summary.png"):
    add_rect(s, Inches(5.95), Inches(3.45), Inches(3.75), Inches(1.65),
             fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    s.shapes.add_picture(f"{ASSETS}/04_telegram_summary.png",
                         Inches(7.40), Inches(3.50), width=Inches(1.30),
                         height=Inches(1.55))
add_text(s, Inches(5.95), Inches(5.13), Inches(3.75), Inches(0.20),
         "Telegram bot — notify_complete card",
         size=9, color=BLUE_LIGHT)



# ── Slide 6 — Methodology / Implementation Details ──────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Methodology — Implementation Details",
                 "How each layer is built and how false positives are suppressed")
add_footer(s)

def box(slide, x, y, w, h, title, lines, accent):
    add_rect(slide, x, y, w, h, fill=BLUE_PALE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.05)
    add_rect(slide, x, y, w, Inches(0.35), fill=accent,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.15)
    add_text(slide, x + Inches(0.15), y + Inches(0.03), w - Inches(0.2),
             Inches(0.30), title, size=11.5, bold=True, color=WHITE)
    # lines: list of (label, text)
    inner_y = y + Inches(0.45)
    for lbl, txt in lines:
        add_text(slide, x + Inches(0.15), inner_y, Inches(1.5), Inches(0.30),
                 lbl, size=9.5, bold=True, color=NAVY)
        add_text(slide, x + Inches(1.30), inner_y, w - Inches(1.4),
                 Inches(0.30), txt, size=9.5, color=SLATE)
        inner_y += Inches(0.32)

# 2x2 grid of method boxes
box(s, Inches(0.40), Inches(1.10), Inches(4.55), Inches(1.85),
    "Recon & Liveness",
    [
        ("Subdomains:", "amass + subfinder + crt.sh; deduped, in-scope filtered"),
        ("Liveness:",   "DNS / ICMP / HTTP triage — kept only when ≥1 responds"),
        ("Storage:",    "per-module JSON under output/<date>_<target>/"),
    ], BLUE)

box(s, Inches(5.05), Inches(1.10), Inches(4.55), Inches(1.85),
    "Active Probes",
    [
        ("Engines:",   "dalfox (XSS), sqlmap (SQLi), in-house IDOR/SSRF/SSTI/JWT"),
        ("Markers:",   "Per-request UUID markers + baseline comparison"),
        ("Verdict:",   "candidate / confirmed / manual_review per finding"),
    ], BLUE)

box(s, Inches(0.40), Inches(3.05), Inches(4.55), Inches(2.05),
    "False-Positive Guards (this work)",
    [
        ("Cache-poison:",     "drop reflections on no-cache/private responses"),
        ("Host-header:",      "downgrade PWA-manifest / canonical reflections"),
        ("IDOR/BOLA:",        "skip cache-buster params and static-asset paths"),
        ("SSTI / XSS:",       "reject debugbar JSON-URI echo reflections"),
        ("Sensitive files:",  "SPA-fallback baseline + ctype/magic signatures"),
        ("Secrets:",          "GitHub repo relevance scoring per target org"),
    ], RED)

box(s, Inches(5.05), Inches(3.05), Inches(4.55), Inches(2.05),
    "Reporting Pipeline",
    [
        ("HTML:",      "Jinja2 templates with filterable finding tables"),
        ("AI report:", "Ollama or OpenAI-compatible API; deterministic fallback"),
        ("Asset risk:", "Weighted score per host → tiered triage list"),
        ("Telegram:",  "Severity-summary bot post on scan completion"),
        ("Diff mode:", "Compare runs to highlight new/removed findings"),
        ("Stack:",     "Python 3.13 • requests • httpx • Jinja2 • python-pptx"),
    ], NAVY)


# ── Slide 7 — Experimental Results ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Experimental Results",
                 "Live scan against miras.app — before vs. after FP-reduction work")
add_footer(s)

# KPIs row
kpis = [
    ("49",     "Subdomains discovered",   BLUE),
    ("39",     "Live HTTP hosts",         BLUE),
    ("42",     "Open ports",              BLUE),
    ("89,673", "Endpoints crawled",       BLUE),
]
for i, (big, lbl, c) in enumerate(kpis):
    x = Inches(0.40 + i * 2.35)
    add_rect(s, x, Inches(1.10), Inches(2.20), Inches(0.95),
             fill=BLUE_PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.05)
    add_text(s, x + Inches(0.10), Inches(1.15), Inches(2.0), Inches(0.50),
             big, size=22, bold=True, color=c)
    add_text(s, x + Inches(0.10), Inches(1.65), Inches(2.0), Inches(0.40),
             lbl, size=10, color=SLATE)

# FP-reduction table
add_text(s, Inches(0.40), Inches(2.25), Inches(8), Inches(0.3),
         "FALSE-POSITIVE REDUCTION (per category)",
         size=10, bold=True, color=RED)

rows = [
    ("Bucket",                       "Before",   "After",   "Reduction"),
    ("Cache poisoning (HIGH)",       "73",       "~few",    "↓ ~95%"),
    ("Host-header reflection",       "185",      "≈30",     "↓ 84%"),
    ("IDOR/BOLA candidates",         "200 (cap)", "≈30",    "↓ 85%"),
    ("Prototype-pollution (HIGH)",   "43",       "0",       "↓ 100%"),
    ("API endpoints classified",     "88",       "55",      "↓ 38%"),
    ("Auth endpoints classified",    "747",      "125",     "↓ 83%"),
    ("Interesting directories",      "3,300",    "193",     "↓ 94%"),
]
table_x = Inches(0.40)
table_y = Inches(2.55)
col_w = [Inches(3.4), Inches(1.6), Inches(1.6), Inches(2.4)]
row_h = Inches(0.30)
x = table_x
for i, (a, b, c, d) in enumerate(rows):
    is_head = (i == 0)
    cells = [a, b, c, d]
    cell_x = table_x
    for j, val in enumerate(cells):
        fill = NAVY if is_head else (BLUE_PALE if i % 2 else WHITE)
        color = WHITE if is_head else NAVY
        bold = is_head or (j == 0)
        add_rect(s, cell_x, table_y + i * row_h, col_w[j], row_h,
                 fill=fill, line=GRAY_LIGHT)
        add_text(s, cell_x + Inches(0.10), table_y + i * row_h, col_w[j],
                 row_h, val, size=10, bold=bold, color=color,
                 anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        cell_x += col_w[j]

# Bug fixed note
add_rect(s, Inches(0.40), Inches(4.95), Inches(9.20), Inches(0.30),
         fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.20)
add_text(s, Inches(0.55), Inches(4.97), Inches(9.0), Inches(0.27),
         "Telegram-bot crash fixed:  ai_report._build_prompt now handles int values "
         "in fuzzer.classified (TypeError: object of type 'int' has no len()).",
         size=10, bold=True, color=WHITE)

# small reference screenshot in upper-right corner
import os as _os
if _os.path.exists(f"{ASSETS}/01_top_risk_assets.png"):
    s.shapes.add_picture(f"{ASSETS}/01_top_risk_assets.png",
                         Inches(7.10), Inches(0.30), width=Inches(2.50),
                         height=Inches(0.72))


# ── Slide 8 — Scientific Novelty & Comparative Analysis ──────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Scientific Novelty & Comparative Analysis",
                 "ReconX vs. existing reconnaissance / scanning tools")
add_footer(s)

comp_rows = [
    ("Solution",        "Passive Recon", "Active Probes", "AI Report", "FP Guards", "OSS / Free", "Single Run"),
    ("Nuclei",          "Limited",       "Yes",           "No",        "Templates", "Yes",        "Yes"),
    ("Burp Suite Pro",  "No",            "Yes",           "No",        "Manual",    "No",         "Manual"),
    ("OWASP ZAP",       "Limited",       "Yes",           "No",        "Limited",   "Yes",        "Manual"),
    ("Amass + Subfinder", "Yes",         "No",            "No",        "n/a",       "Yes",        "Partial"),
    ("ReconX (this work)", "Yes",        "Yes",           "Yes",       "Yes",       "Yes",        "Yes"),
]
col_widths = [Inches(2.20), Inches(1.20), Inches(1.20), Inches(1.10),
              Inches(1.20), Inches(1.10), Inches(1.20)]
table_x = Inches(0.30)
table_y = Inches(1.10)
row_h = Inches(0.30)
for i, row in enumerate(comp_rows):
    is_head = (i == 0)
    is_last = (i == len(comp_rows) - 1)
    cell_x = table_x
    for j, val in enumerate(row):
        if is_head:
            fill, color, bold = NAVY, WHITE, True
        elif is_last:
            fill, color, bold = BLUE, WHITE, True
        else:
            fill = BLUE_PALE if i % 2 else WHITE
            color, bold = NAVY, (j == 0)
        add_rect(s, cell_x, table_y + i * row_h, col_widths[j], row_h,
                 fill=fill, line=GRAY_LIGHT)
        # colour cell values for Yes/No
        cell_color = color
        if val == "Yes" and not is_head and not is_last:
            cell_color = GREEN
        elif val == "No":
            cell_color = RED if not is_last and not is_head else color
        elif val in ("Limited", "Partial", "Manual", "Templates"):
            cell_color = AMBER if not is_last and not is_head else color
        add_text(s, cell_x + Inches(0.05), table_y + i * row_h, col_widths[j],
                 row_h, val, size=10, bold=bold, color=cell_color,
                 anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        cell_x += col_widths[j]

# Novelty bullets
y = table_y + len(comp_rows) * row_h + Inches(0.15)
add_text(s, Inches(0.40), y, Inches(9), Inches(0.30),
         "KEY CONTRIBUTIONS (NOVELTY)",
         size=10, bold=True, color=BLUE)
contrib = [
    ("1", "First OSS framework combining passive recon, 20+ active probes, AI-generated client report, and per-module FP guards in a single pipeline."),
    ("2", "Reproducible FP-reduction methodology: every category (cache_poison, IDOR, host-header, SSTI, XSS, sensitive_files, secrets) has a verifier described in this work."),
    ("3", "Repo-relevance scoring filters unrelated public-GitHub secrets before they pollute analyst reports — addresses a noise source not handled by Gitleaks or trufflehog."),
]
for i, (n, t) in enumerate(contrib):
    yy = y + Inches(0.30 + i * 0.40)
    add_rect(s, Inches(0.40), yy, Inches(0.30), Inches(0.30),
             fill=BLUE, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(0.40), yy, Inches(0.30), Inches(0.30),
             n, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.80), yy, Inches(8.7), Inches(0.40),
             t, size=10, color=SLATE)


# ── Slide 9 — Limitations & Future Work ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Limitations & Future Work",
                 "Honest assessment and planned improvements")
add_footer(s)

# Left column — Limitations
add_text(s, Inches(0.40), Inches(1.10), Inches(4.50), Inches(0.30),
         "LIMITATIONS", size=11, bold=True, color=RED)
lims = [
    ("Out-of-scope coverage",   "External recon only; internal-network pivot and post-auth scanning are out of scope."),
    ("LLM dependency",          "AI report quality varies by model. Ollama deepseek-r1:7b works; weaker local models hallucinate."),
    ("Active-probe rate limits","Probes throttle conservatively to avoid DoS — full-coverage scans on large surfaces still take hours."),
    ("Authenticated flows",     "Sessions are supported via cookie profiles but multi-step login automation needs manual scripting."),
    ("Verifier coverage",       "FP guards cover the most-noisy categories; long-tail rules (e.g. some Nuclei templates) still need triage."),
]
for i, (h, body) in enumerate(lims):
    y = Inches(1.45 + i * 0.66)
    add_rect(s, Inches(0.40), y, Inches(0.10), Inches(0.55), fill=RED)
    add_text(s, Inches(0.60), y, Inches(4.40), Inches(0.30),
             h, size=11, bold=True, color=NAVY)
    add_text(s, Inches(0.60), y + Inches(0.25), Inches(4.40), Inches(0.40),
             body, size=9.5, color=SLATE)

# Right column — Future work
add_text(s, Inches(5.10), Inches(1.10), Inches(4.50), Inches(0.30),
         "FUTURE WORK", size=11, bold=True, color=GREEN)
future = [
    "Authenticated-scan profiles: integrate Playwright for multi-step login replay.",
    "Distributed scanning: shard live-host probes across worker nodes.",
    "Vector-database memory: store findings across runs for diff and trend analysis.",
    "Fine-tuned local LLM for security reporting (RAG over OWASP / CVE corpora).",
    "Continuous-scan mode: hourly diff scans against the same surface, alert on deltas.",
    "Browser-based dashboard: replace static HTML with a live FastAPI + React UI.",
]
for i, line in enumerate(future):
    y = Inches(1.45 + i * 0.55)
    add_text(s, Inches(5.20), y, Inches(0.30), Inches(0.30),
             f"{i+1}.", size=11, bold=True, color=GREEN)
    add_text(s, Inches(5.55), y, Inches(4.10), Inches(0.55),
             line, size=10, color=SLATE)


# ── Slide — Tools & Technology Stack ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "Tools & Technology Stack",
                 "What ReconX is built with and what it integrates")
add_footer(s)

def stack_card(slide, x, y, w, h, title, items, head_color):
    add_rect(slide, x, y, w, h, fill=BLUE_PALE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    add_rect(slide, x, y, w, Inches(0.32), fill=head_color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.12)
    add_text(slide, x + Inches(0.15), y + Inches(0.03), w - Inches(0.2),
             Inches(0.30), title, size=11, bold=True, color=WHITE)
    inner_y = y + Inches(0.40)
    for item in items:
        add_text(slide, x + Inches(0.15), inner_y, w - Inches(0.30),
                 Inches(0.22), "• " + item, size=9.5, color=SLATE)
        inner_y += Inches(0.22)

# 3 columns × 2 rows of cards
stack_card(s, Inches(0.30), Inches(1.10), Inches(3.10), Inches(1.95),
           "Core / Language",
           ["Python 3.13 (CLI: main.py)",
            "Concurrent recon via ThreadPoolExecutor",
            "requests + urllib3 (HTTP)",
            "httpx (async probes)",
            "Jinja2 (HTML templating)",
            "python-pptx (this deck)"],
           BLUE)

stack_card(s, Inches(3.50), Inches(1.10), Inches(3.10), Inches(1.95),
           "Recon / Discovery Tools",
           ["amass · subfinder · assetfinder",
            "crt.sh certificate transparency",
            "nmap / masscan / rustscan",
            "httpx · gowitness (screenshots)",
            "wappalyzer / WhatWeb (techstack)",
            "testssl.sh (TLS audit)"],
           NAVY)

stack_card(s, Inches(6.70), Inches(1.10), Inches(3.00), Inches(1.95),
           "Active Probes / Scanners",
           ["nuclei (template engine)",
            "dalfox (XSS)",
            "sqlmap (SQL injection)",
            "gitleaks · trufflehog (secrets)",
            "ffuf (fuzzer / directory brute)",
            "Interactsh (OOB SSRF callback)"],
           BLUE_BRIGHT)

stack_card(s, Inches(0.30), Inches(3.15), Inches(3.10), Inches(1.90),
           "AI / ML",
           ["Ollama (local LLM runtime)",
            "deepseek-r1 · qwen2.5 · llama3",
            "OpenAI-compatible API fallback",
            "markdown → HTML renderer",
            "bleach (HTML sanitiser)",
            "Deterministic static fallback"],
           BLUE)

stack_card(s, Inches(3.50), Inches(3.15), Inches(3.10), Inches(1.90),
           "Integration / Output",
           ["python-telegram-bot",
            "JSON + HTML + Markdown reports",
            "Per-module artifact directories",
            "Diff-mode (run-to-run delta)",
            "PyYAML (config.yaml)",
            "dotenv (.env secrets)"],
           NAVY)

stack_card(s, Inches(6.70), Inches(3.15), Inches(3.00), Inches(1.90),
           "Dev / Testing",
           ["pytest (171 tests passing)",
            "pytest-mock · responses",
            "Git + GitHub for VCS",
            "Kali Linux 2026 dev environment",
            "Chromium (headless screenshots)",
            "VS Code / Vim"],
           BLUE_BRIGHT)


# ── Slide — References ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_slide_header(s, "References",
                 "Standards, frameworks, and published research consulted")
add_footer(s)

refs = [
    ("OWASP Top 10 — 2021",
     "owasp.org/Top10/  — A01 Broken Access Control … A10 SSRF. ReconX findings tag CWE/OWASP categories."),
    ("OWASP API Security Top 10 — 2023",
     "owasp.org/API-Security/  — used as the taxonomy for idor_probe and api_schema_audit."),
    ("OWASP ASVS v4 / WSTG",
     "owasp.org/www-project-application-security-verification-standard/ — verification levels for the report."),
    ("PortSwigger Web Security Academy",
     "portswigger.net/web-security — cache poisoning, host header injection, prototype pollution, SSTI references."),
    ("Nuclei — ProjectDiscovery",
     "github.com/projectdiscovery/nuclei — template engine integrated by modules/vulnscan.py."),
    ("Dalfox — Cross-Site-Scripting scanner",
     "github.com/hahwul/dalfox — used by modules/xss.py for confirmed XSS detection."),
    ("Gitleaks — open-source secret scanning",
     "github.com/gitleaks/gitleaks — used by modules/secret_scanner.py."),
    ("MITRE CWE & ATT&CK",
     "cwe.mitre.org · attack.mitre.org — finding metadata tagged with CWE-IDs and ATT&CK techniques."),
    ("NIST SP 800-115 — Technical Guide to Information Security Testing",
     "csrc.nist.gov/publications/detail/sp/800-115/final — methodology baseline for external pentesting."),
    ("CIS Controls v8",
     "cisecurity.org/controls — control mapping for the remediation recommendations."),
    ("Anthropic / Ollama / DeepSeek documentation",
     "ollama.com/library · deepseek.com — LLM runtime and model used by modules/ai_report.py."),
    ("Verizon Data Breach Investigations Report (DBIR) 2024",
     "verizon.com/business/resources/reports/dbir/ — basis for the motivation slide statistic."),
]
y0 = Inches(1.10)
for i, (title, body) in enumerate(refs):
    yy = y0 + Inches(i * 0.33)
    add_text(s, Inches(0.40), yy, Inches(0.32), Inches(0.30),
             f"[{i+1}]", size=10, bold=True, color=BLUE)
    add_text(s, Inches(0.80), yy, Inches(3.30), Inches(0.30),
             title, size=10, bold=True, color=NAVY)
    add_text(s, Inches(4.15), yy, Inches(5.45), Inches(0.30),
             body, size=9.5, color=SLATE)


# ── Slide 10 — Conclusion ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY_DEEP)
add_text(s, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
         "Conclusion", size=30, bold=True, color=WHITE)

cards = [
    ("ReconX Framework",
     "Modular Python pipeline integrating passive recon, 20+ active probes, "
     "and AI-augmented reporting in a single open-source tool."),
    ("False-Positive Reduction",
     "Per-module verifiers cut report noise across cache-poison, host-header, "
     "IDOR, prototype-pollution, SSTI, XSS, sensitive files, and secrets categories."),
    ("Research Gap Closed",
     "First OSS framework satisfying recon + active probes + AI report + FP guards + "
     "single-run delivery — validated on real external targets."),
    ("Practical Output",
     "Filterable HTML report, structured JSON, Telegram summary; reproducible "
     "artifacts under output/<date>_<target>/ enable diff and re-run."),
]
for i, (h, t) in enumerate(cards):
    x = Inches(0.40 + (i % 2) * 4.65)
    y = Inches(1.20 + (i // 2) * 1.55)
    add_rect(s, x, y, Inches(4.55), Inches(1.40),
             fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.05)
    add_rect(s, x, y, Inches(0.10), Inches(1.40), fill=BLUE_BRIGHT)
    add_text(s, x + Inches(0.25), y + Inches(0.10), Inches(4.20),
             Inches(0.40), h, size=14, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.50), Inches(4.20),
             Inches(0.85), t, size=10.5, color=BLUE_LIGHT)

# Thank you band
add_rect(s, Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.60),
         fill=BLUE_BRIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.10)
add_text(s, Inches(0.40), Inches(4.55), Inches(9.20), Inches(0.30),
         "THANK YOU FOR YOUR ATTENTION",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.40), Inches(4.83), Inches(9.20), Inches(0.30),
         f"Questions?   |   Mavlanov Shakhrukh   |   CS-2304   |   Supervisor: {SUPERVISOR}",
         size=10, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/home/kali/Desktop/Diploma-pre-defend/reconx/ReconX_Defense_Mavlanov.pptx"
prs.save(out_path)
print(f"Wrote {out_path}")
